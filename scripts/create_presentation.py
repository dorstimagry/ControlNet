#!/usr/bin/env python3
"""Convert presentation_content.md to PowerPoint (.pptx) format.

This script parses the markdown presentation content and creates a PowerPoint
presentation with proper formatting, slides, and structure.

Usage:
    python scripts/create_presentation.py --input presentation_content.md --output presentation.pptx
"""

from __future__ import annotations

import argparse
import re
import tempfile
from pathlib import Path
from typing import List, Tuple, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    exit(1)

try:
    import matplotlib
    matplotlib.use('Agg')  # Non-interactive backend
    import matplotlib.pyplot as plt
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False
    print("Warning: matplotlib not available. Equations will be rendered as text.")


def parse_markdown_sections(content: str) -> List[dict]:
    """Parse markdown content into sections.
    
    Returns:
        List of section dictionaries with keys: title, level, content, subsections
    """
    sections = []
    lines = content.split('\n')
    
    current_section = None
    current_content = []
    in_code_block = False
    code_block_type = None
    
    for line in lines:
        # Handle code blocks
        if line.strip().startswith('```'):
            if not in_code_block:
                in_code_block = True
                code_block_type = line.strip()[3:].strip()
                if current_section:
                    current_content.append(line)
            else:
                in_code_block = False
                code_block_type = None
                if current_section:
                    current_content.append(line)
            continue
        
        if in_code_block:
            if current_section:
                current_content.append(line)
            continue
        
        # Check for horizontal rule (section separator)
        if line.strip() == '---':
            if current_section and current_content:
                current_section['content'] = '\n'.join(current_content)
                sections.append(current_section)
            current_section = None
            current_content = []
            continue
        
        # Check for headings
        if line.startswith('#'):
            # Save previous section
            if current_section:
                current_section['content'] = '\n'.join(current_content)
                sections.append(current_section)
            
            # Determine heading level
            level = len(line) - len(line.lstrip('#'))
            title = line.lstrip('#').strip()
            
            # Start new section
            current_section = {
                'title': title,
                'level': level,
                'content': '',
                'subsections': []
            }
            current_content = []
        else:
            if current_section:
                current_content.append(line)
    
    # Save last section
    if current_section:
        current_section['content'] = '\n'.join(current_content)
        sections.append(current_section)
    
    return sections


def clean_text(text: str) -> str:
    """Clean markdown formatting from text."""
    # Remove markdown links [text](url) -> text
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    # Remove markdown bold **text** -> text
    text = re.sub(r'\*\*([^\*]+)\*\*', r'\1', text)
    # Remove markdown italic *text* -> text
    text = re.sub(r'\*([^\*]+)\*', r'\1', text)
    # Remove markdown code `code` -> code
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Remove markdown headers
    text = re.sub(r'^#+\s+', '', text, flags=re.MULTILINE)
    return text.strip()


def format_bullet_points(text: str) -> List[str]:
    """Extract bullet points from text."""
    bullets = []
    lines = text.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # Check for markdown bullet points
        if line.startswith('- ') or line.startswith('* '):
            bullet = clean_text(line[2:])
            if bullet:
                bullets.append(bullet)
        # Check for numbered lists
        elif re.match(r'^\d+\.\s+', line):
            bullet = clean_text(re.sub(r'^\d+\.\s+', '', line))
            if bullet:
                bullets.append(bullet)
        # Check for checkboxes
        elif line.startswith('- [x]') or line.startswith('- [X]'):
            bullet = clean_text(line[5:])
            if bullet:
                bullets.append('✓ ' + bullet)
        elif line.startswith('- [ ]'):
            bullet = clean_text(line[5:])
            if bullet:
                bullets.append('○ ' + bullet)
    
    return bullets


def render_latex_formula(formula: str, output_path: Path, dpi: int = 150) -> Optional[Path]:
    """Render a LaTeX formula to an image file.
    
    Args:
        formula: LaTeX formula string
        output_path: Path to save the image
        dpi: Resolution for the image
        
    Returns:
        Path to the saved image, or None if failed
    """
    if not HAS_MATPLOTLIB:
        return None
    
    try:
        # Clean formula - remove extra whitespace
        formula = formula.strip()
        
        fig = plt.figure(figsize=(10, 2))
        fig.patch.set_facecolor('white')
        fig.patch.set_alpha(0)
        
        # Try to render with matplotlib's math text renderer
        # This handles basic LaTeX without needing a full LaTeX installation
        plt.text(0.5, 0.5, f'${formula}$', 
                fontsize=24, 
                ha='center', va='center',
                usetex=False,  # Use matplotlib's built-in renderer
                math_fontfamily='dejavusans')  # Use DejaVu Sans for math symbols
        
        plt.xlim(0, 1)
        plt.ylim(0, 1)
        plt.axis('off')
        plt.tight_layout(pad=0.2)
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        fig.savefig(output_path, dpi=dpi, bbox_inches='tight', 
                   facecolor='white', transparent=False, 
                   pad_inches=0.1)
        plt.close(fig)
        
        return output_path
    except Exception as e:
        print(f"Warning: Could not render formula '{formula[:50]}...': {e}")
        return None


def find_matching_image(figure_description: str, search_dirs: List[Path]) -> Optional[Path]:
    """Find an image file that matches the figure description.
    
    Args:
        figure_description: Description from [FIGURE: ...] placeholder
        search_dirs: List of directories to search for images
        
    Returns:
        Path to matching image file, or None if not found
    """
    # Extract keywords from description
    keywords = re.findall(r'\b\w+\b', figure_description.lower())
    
    # Common image file patterns to look for
    image_extensions = ['.png', '.jpg', '.jpeg', '.pdf']
    
    # Search in directories
    for search_dir in search_dirs:
        if not search_dir.exists():
            continue
        
        # Look for common plot filenames
        for ext in image_extensions:
            # Try exact matches first
            for pattern in ['step_function_results', 'profile_statistics', 
                           'closed_loop_summary', 'episode_000', 'episode_001']:
                candidate = search_dir / f"{pattern}{ext}"
                if candidate.exists():
                    # Check if keywords match
                    if any(kw in pattern.lower() for kw in keywords[:3]):
                        return candidate
            
            # Search recursively for files matching keywords
            for img_path in search_dir.rglob(f'*{ext}'):
                img_name_lower = img_path.stem.lower()
                # Check if any keyword matches
                if any(kw in img_name_lower for kw in keywords[:3] if len(kw) > 3):
                    return img_path
    
    return None


def extract_formulas(text: str) -> Tuple[str, List[str]]:
    """Extract LaTeX formulas from text and replace with placeholders.
    
    Returns:
        (text_with_placeholders, list_of_formulas)
    """
    formulas = []
    # Match \[ ... \] (display math)
    pattern = r'\\\[(.*?)\\\]'
    matches = list(re.finditer(pattern, text, re.DOTALL))
    
    for i, match in enumerate(matches):
        formula = match.group(1).strip()
        formulas.append(formula)
        text = text.replace(match.group(0), f'[FORMULA_{i}]')
    
    # Match \( ... \) (inline math)
    pattern = r'\\\((.*?)\\\)'
    matches = list(re.finditer(pattern, text))
    
    for i, match in enumerate(matches):
        formula = match.group(1).strip()
        formulas.append(formula)
        text = text.replace(match.group(0), f'[FORMULA_{len(formulas)-1}]')
    
    return text, formulas


def extract_figure_placeholders(text: str) -> List[str]:
    """Extract [FIGURE: ...] placeholders from text."""
    pattern = r'\[FIGURE:\s*([^\]]+)\]'
    matches = re.findall(pattern, text, re.IGNORECASE)
    return matches


def create_title_slide(prs: Presentation, title: str, subtitle: str = None):
    """Create title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle


def create_content_slide(prs: Presentation, title: str, content: str, 
                        is_section_header: bool = False,
                        figure_search_dirs: List[Path] = None,
                        formula_temp_dir: Path = None):
    """Create a content slide with title and content.
    
    Args:
        prs: Presentation object
        title: Slide title
        content: Slide content
        is_section_header: Whether this is a section header slide
        figure_search_dirs: Directories to search for figure images
        formula_temp_dir: Temporary directory for formula images
    """
    slide_layout = prs.slide_layouts[1]  # Title and content
    
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_clean = clean_text(title)
    # Remove markdown formatting from title
    title_clean = re.sub(r'\*\*([^\*]+)\*\*', r'\1', title_clean)  # Bold
    title_shape.text = title_clean
    
    # Get content placeholder
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    if not content.strip():
        return
    
    # Extract formulas
    content_clean, formulas = extract_formulas(content)
    
    # Extract figure placeholders
    figures = extract_figure_placeholders(content_clean)
    
    # Try to find and insert actual images
    inserted_images = []
    if figure_search_dirs:
        for fig_desc in figures:
            img_path = find_matching_image(fig_desc, figure_search_dirs)
            if img_path and img_path.exists():
                inserted_images.append((fig_desc, img_path))
                # Remove placeholder from text
                content_clean = re.sub(
                    rf'\[FIGURE:\s*{re.escape(fig_desc)}\]', 
                    '', content_clean, flags=re.IGNORECASE
                )
    
    # Remove remaining figure placeholders
    content_clean = re.sub(r'\[FIGURE:[^\]]+\]', '', content_clean, flags=re.IGNORECASE)
    
    # Split into lines for better processing
    lines = content_clean.split('\n')
    
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        
        if not line:
            i += 1
            continue
        
        # Handle code blocks
        if line.startswith('```'):
            # Collect code block
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            if i < len(lines):
                i += 1
            
            # Add code block as monospace text
            code_text = '\n'.join(code_lines)
            p = text_frame.add_paragraph()
            p.text = code_text
            p.font.name = 'Courier New'
            p.font.size = Pt(9)
            p.space_after = Pt(12)
            continue
        
        # Handle bullet points
        if line.startswith('- ') or line.startswith('* ') or re.match(r'^\d+\.\s+', line):
            bullet_text = clean_text(line)
            # Remove bullet marker
            if bullet_text.startswith('- ') or bullet_text.startswith('* '):
                bullet_text = bullet_text[2:]
            elif re.match(r'^\d+\.\s+', bullet_text):
                bullet_text = re.sub(r'^\d+\.\s+', '', bullet_text)
            
            p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0
            p.font.size = Pt(14)
            p.space_after = Pt(6)
            i += 1
            continue
        
        # Handle blockquotes
        if line.startswith('>'):
            quote_text = clean_text(line[1:].strip())
            p = text_frame.add_paragraph()
            p.text = quote_text
            p.font.italic = True
            p.level = 0
            p.left_indent = Inches(0.5)
            i += 1
            continue
        
        # Handle tables (simple detection)
        if '|' in line and line.count('|') >= 2:
            # Collect table rows
            table_lines = []
            while i < len(lines) and '|' in lines[i]:
                table_lines.append(lines[i])
                i += 1
            
            # Add table as formatted text
            for table_line in table_lines:
                # Clean table formatting
                cells = [c.strip() for c in table_line.split('|') if c.strip()]
                if cells and not all(c.startswith('-') for c in cells):  # Skip separator row
                    table_text = ' | '.join(cells)
                    p = text_frame.add_paragraph()
                    p.text = table_text
                    p.font.size = Pt(11)
                    p.font.name = 'Courier New'
            continue
        
        # Regular paragraph - collect until blank line
        para_lines = []
        while i < len(lines) and lines[i].strip():
            para_lines.append(lines[i])
            i += 1
        
        if para_lines:
            para_text = ' '.join(para_lines)
            para_clean = clean_text(para_text)
            
            # Replace formula placeholders
            for idx, formula in enumerate(formulas):
                para_clean = para_clean.replace(f'[FORMULA_{idx}]', f'[{formula}]')
            
            if para_clean:
                p = text_frame.add_paragraph()
                p.text = para_clean
                p.level = 0
                p.font.size = Pt(14)
                p.space_after = Pt(12)
    
    # Insert images if found
    if inserted_images:
        # Calculate image positions
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(4.5)
        
        for fig_desc, img_path in inserted_images:
            try:
                slide.shapes.add_picture(str(img_path), left, top, width, height)
                # Stack images vertically if multiple
                top += height + Inches(0.2)
            except Exception as e:
                print(f"Warning: Could not insert image {img_path}: {e}")
    
    # Render and insert formula images
    if formulas and formula_temp_dir:
        formula_images = []
        for i, formula in enumerate(formulas):
            formula_path = formula_temp_dir / f"formula_{i}.png"
            rendered = render_latex_formula(formula, formula_path)
            if rendered and rendered.exists():
                formula_images.append((i, rendered))
        
        # Insert formula images
        if formula_images:
            formula_left = Inches(1)
            formula_top = Inches(4.5)
            formula_width = Inches(8)
            
            for i, formula_path in formula_images:
                try:
                    # Calculate height based on image
                    try:
                        from PIL import Image
                        img = Image.open(formula_path)
                        img_width, img_height = img.size
                        aspect_ratio = img_height / img_width
                        formula_height = formula_width * aspect_ratio
                    except ImportError:
                        # Fallback: use fixed height
                        formula_height = Inches(0.8)
                    
                    slide.shapes.add_picture(
                        str(formula_path), 
                        formula_left, formula_top, 
                        formula_width, formula_height
                    )
                    formula_top += formula_height + Inches(0.2)
                except Exception as e:
                    print(f"Warning: Could not insert formula image {formula_path}: {e}")
    
    # Add figure placeholders as notes (for ones not found)
    missing_figures = [fig for fig in figures if fig not in [f[0] for f in inserted_images]]
    if missing_figures:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "Figure Placeholders (not found):\n" + "\n".join(f"- {fig}" for fig in missing_figures)


def create_section_header_slide(prs: Presentation, title: str):
    """Create a section header slide."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title text box
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = clean_text(title)
    text_frame.word_wrap = True
    
    # Format title
    p = text_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def convert_markdown_to_pptx(markdown_path: Path, output_path: Path, 
                            figure_search_dirs: List[Path] = None):
    """Convert markdown presentation to PowerPoint.
    
    Args:
        markdown_path: Path to markdown file
        output_path: Path to output PowerPoint file
        figure_search_dirs: Directories to search for figure images
    """
    print(f"Reading markdown from: {markdown_path}")
    with open(markdown_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    print("Parsing markdown content...")
    sections = parse_markdown_sections(content)
    
    print(f"Found {len(sections)} sections")
    
    # Set up figure search directories
    if figure_search_dirs is None:
        # Default: search in evaluation/results/plots
        project_root = markdown_path.parent
        figure_search_dirs = [
            project_root / "evaluation" / "results" / "plots",
            project_root / "evaluation" / "results" / "plots" / "closed_loop_cpp_comparison_with_initial_error",
        ]
    
    # Create temporary directory for formula images
    with tempfile.TemporaryDirectory() as temp_dir:
        formula_temp_dir = Path(temp_dir)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Extract title slide info
        title_slide_created = False
        for section in sections[:3]:  # Check first few sections
            if section['level'] == 1 and ('Deep RL' in section['title'] or 'SAC' in section['title']):
                # Extract subtitle from content
                subtitle_lines = section['content'].split('\n')[:5]
                subtitle = None
                for line in subtitle_lines:
                    if 'Presenter' in line or 'Date' in line:
                        continue
                    if line.strip() and not line.startswith('**'):
                        subtitle = line.strip()
                        break
                
                create_title_slide(prs, section['title'], subtitle or "Soft Actor-Critic (SAC) Approach")
                title_slide_created = True
                break
        
        # Process sections
        for i, section in enumerate(sections):
            title = section['title']
            content_text = section['content']
            level = section['level']
            
            # Skip title slide section if already created
            if title_slide_created and i == 0 and level == 1:
                continue
            
            print(f"Processing section {i+1}/{len(sections)} (level {level}): {title[:50]}...")
            
            # Handle section headers (level 1 with numbers)
            if level == 1 and re.match(r'^\d+\.', title):
                create_section_header_slide(prs, title)
            
            # Split content if it's too long (rough heuristic)
            if content_text.strip():
                # Split by double newlines to get logical chunks
                chunks = [c.strip() for c in content_text.split('\n\n') if c.strip()]
                
                # If content is very long, split into multiple slides
                if len(content_text) > 2000:
                    # Create first slide with title
                    first_chunk = '\n\n'.join(chunks[:3])  # First 3 chunks
                    create_content_slide(prs, title, first_chunk, 
                                       figure_search_dirs=figure_search_dirs,
                                       formula_temp_dir=formula_temp_dir)
                    
                    # Create continuation slides
                    remaining_chunks = chunks[3:]
                    for j in range(0, len(remaining_chunks), 3):
                        chunk = '\n\n'.join(remaining_chunks[j:j+3])
                        create_content_slide(prs, f"{title} (continued)", chunk,
                                           figure_search_dirs=figure_search_dirs,
                                           formula_temp_dir=formula_temp_dir)
                else:
                    create_content_slide(prs, title, content_text,
                                       figure_search_dirs=figure_search_dirs,
                                       formula_temp_dir=formula_temp_dir)
        
        # Save presentation
        print(f"Saving PowerPoint to: {output_path}")
        prs.save(str(output_path))
        print(f"Done! Created {len(prs.slides)} slides.")


def main():
    parser = argparse.ArgumentParser(description="Convert markdown presentation to PowerPoint")
    parser.add_argument(
        '--input',
        type=Path,
        default=Path('docs/presentation_content.md'),
        help='Input markdown file (default: docs/presentation_content.md)'
    )
    parser.add_argument(
        '--output',
        type=Path,
        default=Path('docs/presentation.pptx'),
        help='Output PowerPoint file (default: docs/presentation.pptx)'
    )
    parser.add_argument(
        '--figure-dirs',
        type=Path,
        nargs='+',
        help='Additional directories to search for figure images'
    )
    
    args = parser.parse_args()
    
    if not args.input.exists():
        print(f"Error: Input file not found: {args.input}")
        return 1
    
    figure_dirs = args.figure_dirs if args.figure_dirs else None
    convert_markdown_to_pptx(args.input, args.output, figure_search_dirs=figure_dirs)
    return 0


if __name__ == '__main__':
    exit(main())

